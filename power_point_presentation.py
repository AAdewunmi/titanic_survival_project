from pptx import Presentation


def create_presentation():
    """
    Creates a PowerPoint presentation that summarizes the Titanic survival
    prediction project.

    The presentation includes:
    - Project overview
    - Model accuracy
    - Classification report
    - Conclusion

    Returns
    -------
    pptx_file : str
        Path to the saved PowerPoint file.
    
    Example Usage:
    --------------
    Create and save the PowerPoint presentation for Titanic Survival Prediction 
    project.
    
    >>> pptx_file = create_presentation()
    >>> print(f"Presentation saved at {pptx_file}")
    """
    # Create a new PowerPoint presentation
    prs = Presentation()

    # Slide 1: Title Slide
    slide_1 = prs.slides.add_slide(
        prs.slide_layouts[0]
    )  # 0: Title slide layout
    title = slide_1.shapes.title  # Access the title of the slide
    subtitle = slide_1.placeholders[1]  # Access the subtitle placeholder

    # Set the title and subtitle text
    title.text = "Titanic Survival Prediction Project"
    subtitle.text = "An analysis and prediction using Logistic Regression"

    # Slide 2: Project Overview
    slide_2 = prs.slides.add_slide(
        prs.slide_layouts[1]
    )  # 1: Title and Content layout
    title = slide_2.shapes.title
    content = slide_2.shapes.placeholders[1]

    # Set project overview content
    title.text = "Project Overview"
    content.text = (
        "This project aims to predict Titanic passenger survival using "
        "machine learning.\n"
        "We use the Kaggle Titanic dataset and clean the data, followed by "
        "building a logistic regression model.\n"
        "The project includes the following steps:\n"
        "1. Data Loading and Cleaning\n"
        "2. Feature Engineering\n"
        "3. Model Training\n"
        "4. Model Evaluation"
    )

    # Slide 3: Accuracy Results
    slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_3.shapes.title
    content = slide_3.shapes.placeholders[1]

    # Set accuracy results content
    title.text = "Model Accuracy"
    content.text = (
        "The logistic regression model achieved an accuracy of 0.8101 on the "
        "test set.\n"
        "This means the model correctly predicted survival for about 81% of "
        "passengers."
    )

    # Slide 4: Classification Report
    slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_4.shapes.title
    content = slide_4.shapes.placeholders[1]

    # Set classification report content
    title.text = "Classification Report"
    content.text = (
        "Here is the detailed classification report for the Titanic survival "
        "prediction model:\n\n"
        "Precision, Recall, F1-Score, and Support for each class "
        "(Survived=0 and Survived=1):\n\n"
        "       precision    recall   f1-score   support\n"
        "   0    0.83        0.86      0.84      105\n"
        "   1    0.79        0.74      0.76      74\n\n"
        "Accuracy: 0.81 (Overall accuracy of the model)\n\n"
        "Macro avg: 0.81  0.80  0.80  179\n"
        "Weighted avg: 0.81  0.81  0.81  179"
    )

    # Slide 5: Conclusion
    slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_5.shapes.title
    content = slide_5.shapes.placeholders[1]

    # Set conclusion content
    title.text = "Conclusion"
    content.text = (
        "The logistic regression model performed well, with an accuracy of "
        "81%.\n"
        "Precision and recall indicate the model is relatively balanced in "
        "predicting survival and non-survival.\n"
        "Future improvements could include exploring more complex models, "
        "feature engineering, and hyperparameter tuning."
    )

    # Save the PowerPoint presentation to a file
    pptx_file = "Titanic_Survival_Prediction_Project.pptx"
    prs.save(pptx_file)

    # Return the path to the saved PowerPoint file
    return pptx_file

# Example usage
if __name__ == "__main__":
    pptx_file = create_presentation()
    print(f"Presentation saved at {pptx_file}")
